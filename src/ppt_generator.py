import json
from pptx import Presentation

# Load insights
with open("../reports/insights.json", "r") as f:
    insights = json.load(f)

# Create presentation
prs = Presentation()

# ---------- Slide 1: Title Slide ----------
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Power AI"
slide.placeholders[1].text = "Automated Data Analysis & Reporting"

# ---------- Slide 2: Dataset Overview ----------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Dataset Overview"
content = slide.placeholders[1]
content.text = (
    f"Total Records: {insights['total_records']}\n"
    f"Total Columns: {insights['total_columns']}"
)

# ---------- Slide 3: Average Salary by Department ----------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Average Salary by Department"
text_frame = slide.placeholders[1].text_frame
text_frame.text = ""

for dept, salary in insights["average_salary_by_department"].items():
    p = text_frame.add_paragraph()
    p.text = f"{dept}: {salary}"
    p.level = 1

# ---------- Slide 4: Key Insight ----------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Key Insight"
slide.placeholders[1].text = (
    f"Highest salary is in the "
    f"{insights['highest_salary']['department']} department "
    f"with a value of {insights['highest_salary']['value']}."
)

# Save presentation
prs.save("../reports/PowerAI_Presentation.pptx")

print("✅ PowerAI_Presentation.pptx created successfully")
